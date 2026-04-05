"""
make_deck.py — Generates the Red Hat Summit 2026 demo slide deck.
Run with: uv run --with python-pptx python slides/make_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Theme ──────────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)
TEAL       = RGBColor(0x00, 0xC2, 0xA8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY       = RGBColor(0xB0, 0xBC, 0xCC)
CARD       = RGBColor(0x14, 0x2A, 0x3E)
RED        = RGBColor(0xEF, 0x44, 0x44)
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
AMBER      = RGBColor(0xF5, 0x9E, 0x0E)

SLIDE_W    = Inches(13.33)
SLIDE_H    = Inches(7.5)

LOGO_PATH  = "/Users/myork/manshed/virtus-slides/netbox_logo_bright.png"
ARCH_PNG   = os.path.join(os.path.dirname(__file__), "architecture.png")
SEQ_PNG    = os.path.join(os.path.dirname(__file__), "sequence.png")

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]  # completely blank


# ── Helpers ────────────────────────────────────────────────────────────────────

def new_slide():
    s = prs.slides.add_slide(blank)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    return s


def add_logo(slide, size=Inches(1.4)):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(
            LOGO_PATH,
            left=SLIDE_W - size - Inches(0.25),
            top=SLIDE_H - size * 0.5 - Inches(0.18),
            width=size,
        )


def txb(slide, text, left, top, width, height,
        size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = "Calibri"
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def rect(slide, left, top, width, height, fill=CARD, line=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def accent_bar(slide, top=Inches(0.08)):
    """Thin teal bar at the top of the slide."""
    r = slide.shapes.add_shape(1, 0, top, SLIDE_W, Inches(0.06))
    r.fill.solid()
    r.fill.fore_color.rgb = TEAL
    r.line.fill.background()


def bullet_slide(slide, title, bullets, logo=True):
    accent_bar(slide)
    txb(slide, title,
        Inches(0.5), Inches(0.25), Inches(11), Inches(0.75),
        size=32, bold=True, color=TEAL)
    y = Inches(1.2)
    for b in bullets:
        indent = b.startswith("  ")
        text   = b.lstrip()
        col    = GREY if indent else WHITE
        sz     = 18 if indent else 22
        prefix = "    •  " if indent else "•  "
        txb(slide, prefix + text,
            Inches(0.6), y, Inches(11.5), Inches(0.55),
            size=sz, color=col)
        y += Inches(0.58) if not indent else Inches(0.5)
    if logo:
        add_logo(slide)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()

# Big teal bar left edge
r = s.shapes.add_shape(1, 0, 0, Inches(0.18), SLIDE_H)
r.fill.solid(); r.fill.fore_color.rgb = TEAL; r.line.fill.background()

txb(s, "Automated Circuit Failover",
    Inches(0.5), Inches(1.6), Inches(9), Inches(1.1),
    size=44, bold=True, color=WHITE)
txb(s, "NetBox Cloud  ·  Ansible Automation Platform",
    Inches(0.5), Inches(2.75), Inches(9), Inches(0.6),
    size=24, color=TEAL)
txb(s, "Red Hat Summit 2026",
    Inches(0.5), Inches(3.5), Inches(6), Inches(0.5),
    size=20, color=GREY)
txb(s, "Matt York  |  Sales Engineer, NetBox Labs",
    Inches(0.5), Inches(4.1), Inches(6), Inches(0.45),
    size=18, color=GREY, italic=True)

add_logo(s, size=Inches(2.0))


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2 — Customer scenario: The Incident
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
accent_bar(s)

txb(s, "A Customer Came to Us With a Problem",
    Inches(0.5), Inches(0.2), Inches(12), Inches(0.75),
    size=32, bold=True, color=TEAL)

# Large quote / scenario box
rect(s, Inches(0.4), Inches(1.1), Inches(12.5), Inches(3.5), fill=CARD, line=RED)

txb(s, "\u201cOur intercontinental IPLC link between Bristol and Atlanta went down.",
    Inches(0.7), Inches(1.3), Inches(12.0), Inches(0.6),
    size=20, bold=True, color=WHITE)
txb(s, "We had a secondary link available. But it took our team 45 minutes to find it,",
    Inches(0.7), Inches(1.95), Inches(12.0), Inches(0.5),
    size=20, color=WHITE)
txb(s, "identify which routers were affected, and manually reconfigure the failover path.",
    Inches(0.7), Inches(2.45), Inches(12.0), Inches(0.5),
    size=20, color=WHITE)
txb(s, "In that 45 minutes, we lost hundreds of thousands of dollars in business.",
    Inches(0.7), Inches(2.95), Inches(12.0), Inches(0.5),
    size=20, color=WHITE)
txb(s, "We need to make sure this never happens again.\u201d",
    Inches(0.7), Inches(3.5), Inches(12.0), Inches(0.5),
    size=20, bold=True, color=RED)

# Attribution
txb(s, "Enterprise Network Customer  |  Intercontinental WAN Operations",
    Inches(0.7), Inches(4.75), Inches(10), Inches(0.4),
    size=14, color=GREY, italic=True)

add_logo(s)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 3 — Why did it take 45 minutes?
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
accent_bar(s)

txb(s, "Why Did It Take 45 Minutes?",
    Inches(0.5), Inches(0.2), Inches(12), Inches(0.75),
    size=32, bold=True, color=TEAL)

timeline = [
    ("0:00", RED,   "Alert fires. Engineers are paged."),
    ("0:05", AMBER, "Team opens spreadsheets and old Visio diagrams to find the secondary circuit."),
    ("0:15", AMBER, "Secondary circuit details found — but which routers does it use?"),
    ("0:22", AMBER, "Cross-reference IPAM, check interface logs, ask a colleague who \u201cknows the network\u201d."),
    ("0:30", AMBER, "SSH into Bristol router. Find the right interface. Apply config."),
    ("0:37", AMBER, "SSH into Atlanta router. Repeat."),
    ("0:43", AMBER, "Test connectivity. Something is wrong. Roll back. Try again."),
    ("0:45", GREEN, "Traffic restored. Documentation updated... maybe."),
]

y = Inches(1.1)
for time, col, desc in timeline:
    rect(s, Inches(0.4), y, Inches(1.2), Inches(0.42), fill=CARD, line=col)
    txb(s, time, Inches(0.4), y + Inches(0.04), Inches(1.2), Inches(0.38),
        size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, desc, Inches(1.75), y + Inches(0.05), Inches(11.0), Inches(0.38),
        size=15, color=WHITE)
    y += Inches(0.52)

add_logo(s)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 4 — The root cause
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
accent_bar(s)

txb(s, "The Root Cause",
    Inches(0.5), Inches(0.2), Inches(12), Inches(0.75),
    size=32, bold=True, color=TEAL)

# Three cards
cards = [
    (AMBER, "No single source of truth",
     "Circuit inventory lived in spreadsheets, Visio diagrams, and people\u2019s heads. No one system had the full picture."),
    (RED,   "No automation",
     "Every failover step was manual: find the backup, identify the routers, SSH in, apply config, test, document."),
    (RED,   "No audit trail",
     "After the incident, the team couldn\u2019t easily answer: what changed, when, and who made the call?"),
]

cx = Inches(0.4)
for col, title, body in cards:
    rect(s, cx, Inches(1.2), Inches(4.0), Inches(5.5), fill=CARD, line=col)
    txb(s, title, cx + Inches(0.2), Inches(1.4), Inches(3.6), Inches(0.7),
        size=20, bold=True, color=col)
    txb(s, body,  cx + Inches(0.2), Inches(2.25), Inches(3.6), Inches(4.0),
        size=17, color=WHITE)
    cx += Inches(4.45)

txb(s, "The technology existed to fix all three of these problems. It just wasn\u2019t joined up.",
    Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.5),
    size=18, color=GREY, italic=True, align=PP_ALIGN.CENTER)

add_logo(s)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 5 — What we recommended
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
accent_bar(s)

txb(s, "What We Recommended",
    Inches(0.5), Inches(0.2), Inches(12), Inches(0.75),
    size=32, bold=True, color=TEAL)

# Left: problem summary box
rect(s, Inches(0.4), Inches(1.1), Inches(4.0), Inches(5.8), fill=CARD, line=RED)
txb(s, "The old world",
    Inches(0.6), Inches(1.25), Inches(3.6), Inches(0.5),
    size=18, bold=True, color=RED)
old = [
    "Spreadsheet inventory",
    "Manual failover runbooks",
    "45+ minute MTTR",
    "No automation",
    "Stale documentation",
    "Errors under pressure",
]
oy = Inches(1.9)
for item in old:
    txb(s, "\u2717  " + item, Inches(0.6), oy, Inches(3.5), Inches(0.45),
        size=16, color=GREY)
    oy += Inches(0.5)

# Arrow
txb(s, "\u2192", Inches(4.55), Inches(3.7), Inches(0.6), Inches(0.6),
    size=40, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Right: solution summary box
rect(s, Inches(5.3), Inches(1.1), Inches(7.6), Inches(5.8), fill=CARD, line=TEAL)
txb(s, "The NetBox + AAP world",
    Inches(5.5), Inches(1.25), Inches(7.2), Inches(0.5),
    size=18, bold=True, color=TEAL)
new_items = [
    ("NetBox Cloud", "Single source of truth for every circuit, device, and connection"),
    ("Copilot AI", "Translate operator intent into precise API actions, no CLI needed"),
    ("Event-driven webhook", "A status change in NetBox triggers automation instantly"),
    ("AAP Controller", "Discovers the backup dynamically, pushes config, updates NetBox"),
    ("Visual Explorer", "Live topology map reflects the new state immediately"),
    ("Full audit trail", "Every change recorded in NetBox and AAP job history"),
]
ny = Inches(1.9)
for label, desc in new_items:
    txb(s, "\u2713  " + label, Inches(5.5), ny, Inches(2.4), Inches(0.45),
        size=16, bold=True, color=TEAL)
    txb(s, desc, Inches(7.9), ny, Inches(4.8), Inches(0.45),
        size=15, color=WHITE)
    ny += Inches(0.5)

add_logo(s)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 6 — The Solution (how it works)
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bullet_slide(s, "How It Works", [
    "NetBox Cloud holds every circuit, device, and connection as structured data",
    "Copilot AI translates plain English into a precise API action",
    "A circuit status change fires a webhook to AAP automatically",
    "AAP discovers the best backup circuit dynamically from NetBox",
    "Router config is pushed and NetBox is updated in a single workflow",
    "Visual Explorer map reflects the new topology in real time",
    "Automated failover: under 30 seconds  \u00b7  Full audit trail",
])


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 7 — Architecture diagram
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
accent_bar(s)
txb(s, "Architecture",
    Inches(0.5), Inches(0.25), Inches(10), Inches(0.7),
    size=32, bold=True, color=TEAL)

if os.path.exists(ARCH_PNG):
    s.shapes.add_picture(ARCH_PNG,
        Inches(0.3), Inches(1.05),
        width=Inches(12.0))

add_logo(s)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 8 — Demo sequence diagram
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
accent_bar(s)
txb(s, "Demo Sequence",
    Inches(0.5), Inches(0.25), Inches(10), Inches(0.7),
    size=32, bold=True, color=TEAL)

if os.path.exists(SEQ_PNG):
    s.shapes.add_picture(SEQ_PNG,
        Inches(0.5), Inches(1.0),
        width=Inches(11.8))

add_logo(s)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 9 — Step-by-step demo flow
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
accent_bar(s)
txb(s, "Demo Flow \u2014 Step by Step",
    Inches(0.5), Inches(0.25), Inches(11), Inches(0.7),
    size=32, bold=True, color=TEAL)

steps = [
    ("1", "Visual Explorer", "Open the global WAN map \u2014 GB-Bristol hub with active circuit to US-Atlanta via IPLC-GB-AT-PRI"),
    ("2", "Copilot AI",      'Tell Copilot: \u201cIPLC-GB-AT-PRI has failed \u2014 set it to offline\u201d'),
    ("3", "AAP Workflow",    "NetBox event rule fires webhook \u2192 AAP launches Circuit Failover Workflow automatically"),
    ("4", "Step 1: Failover","Playbook queries NetBox, finds backup (IPLC-GB-AT-SEC), pushes router config, updates NetBox"),
    ("5", "Step 2: Report",  "Second workflow step generates HTML incident report and deploys it to the report server"),
    ("6", "Visual Explorer", "Refresh map \u2014 primary circuit gone, backup confirmed active on new router path"),
    ("7", "Claude + MCP",    'Ask Claude: \u201cWhat is the current status of IPLC-GB-AT-PRI?\u201d \u2014 live NetBox data confirms'),
]

y = Inches(1.1)
for num, label, desc in steps:
    nr = s.shapes.add_shape(1, Inches(0.35), y + Inches(0.05), Inches(0.38), Inches(0.38))
    nr.fill.solid(); nr.fill.fore_color.rgb = TEAL; nr.line.fill.background()
    txb(s, num, Inches(0.36), y + Inches(0.02), Inches(0.38), Inches(0.38),
        size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txb(s, label, Inches(0.85), y, Inches(2.2), Inches(0.45),
        size=16, bold=True, color=TEAL)
    txb(s, desc,  Inches(3.1), y, Inches(9.8), Inches(0.45),
        size=15, color=WHITE)
    y += Inches(0.52)

add_logo(s)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 10 — What the playbook does (no hardcoding)
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bullet_slide(s, "The Playbook: Driven Entirely by NetBox Data", [
    "Input: just the failed circuit CID (e.g. IPLC-GB-AT-PRI)",
    "Discovers A-side and Z-side sites from circuit terminations in NetBox",
    "Finds all active circuits between those two sites tagged for failover",
    "Selects the best backup by highest committed bandwidth",
    "No hardcoded backup mappings \u2014 add a new circuit to NetBox and it is automatically considered",
    "Pushes router config changes to all routers at both ends",
    "Writes back to NetBox: primary \u2192 offline, backup \u2192 active",
    "Generates and publishes the HTML failover incident report",
])


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 11 — ROI / Key messages
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
accent_bar(s)
txb(s, "The Business Case",
    Inches(0.5), Inches(0.25), Inches(11), Inches(0.7),
    size=32, bold=True, color=TEAL)

# Two metric boxes
for lx, label, val, col in [
    (Inches(0.5),  "Manual failover", "~45 min", RED),
    (Inches(6.8),  "Automated failover", "< 30 sec", GREEN),
]:
    rect(s, lx, Inches(1.1), Inches(5.9), Inches(2.2), fill=CARD, line=col)
    txb(s, val,   lx + Inches(0.2), Inches(1.25), Inches(5.5), Inches(1.1),
        size=52, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, label, lx + Inches(0.2), Inches(2.1), Inches(5.5), Inches(0.5),
        size=18, color=GREY, align=PP_ALIGN.CENTER)

messages = [
    "NetBox is not just a CMDB \u2014 it is an active source of truth that drives automation",
    "The moment a circuit changes state, the response is already running",
    "The playbook discovers backup paths dynamically \u2014 works for any circuit, any site",
    "Every failover is logged in NetBox, captured in AAP job history, and published as a report",
    "MCP lets Claude query live NetBox data \u2014 confirm outcomes in natural language",
]

y = Inches(3.55)
for msg in messages:
    txb(s, "\u2756  " + msg, Inches(0.5), y, Inches(12.3), Inches(0.48),
        size=18, color=WHITE)
    y += Inches(0.52)

add_logo(s)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 12 — MCP: Confirm with Claude (final demo step)
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
accent_bar(s)
txb(s, "Step 7: Confirm with Claude + NetBox MCP",
    Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
    size=32, bold=True, color=TEAL)

# Left panel
rect(s, Inches(0.3), Inches(1.05), Inches(4.5), Inches(5.9), fill=CARD, line=TEAL)
txb(s, "What is the NetBox MCP Server?",
    Inches(0.5), Inches(1.2), Inches(4.1), Inches(0.5),
    size=16, bold=True, color=TEAL)

mcp_bullets = [
    "Runs on a dedicated AWS VM",
    "Connected to your NetBox Cloud instance",
    "Read-only \u2014 queries only, no writes",
    "Accessed via SSH stdio transport",
    "Registered once with Claude Code:\nclaude mcp add netbox-mcp ...",
    "Claude calls NetBox API tools\nautomatically behind the scenes",
]
by = Inches(1.8)
for b in mcp_bullets:
    txb(s, "\u2022  " + b, Inches(0.5), by, Inches(4.1), Inches(0.65),
        size=14, color=WHITE)
    by += Inches(0.62)

# Right panel
rect(s, Inches(5.1), Inches(1.05), Inches(7.9), Inches(5.9), fill=CARD, line=GREY)
txb(s, "Example queries at the end of the demo",
    Inches(5.3), Inches(1.2), Inches(7.5), Inches(0.5),
    size=16, bold=True, color=GREY)

queries = [
    ("What is the current status of IPLC-GB-AT-PRI?",
     "Confirms: Offline \u2014 the failed circuit is recorded"),
    ("What is the status of IPLC-GB-AT-SEC?",
     "Confirms: Active \u2014 the backup is live"),
    ("Show me all active circuits between GB-Bristol and US-Atlanta.",
     "Lists the surviving circuits on that route"),
    ("Which circuits are currently not active?",
     "Full view of any degraded paths across the WAN"),
    ("When was IPLC-GB-AT-PRI last changed, and by whom?",
     "Audit trail: timestamp and username from NetBox change log"),
]

qy = Inches(1.85)
for q, ans in queries:
    qbox = rect(s, Inches(5.2), qy, Inches(7.6), Inches(0.42), fill=RGBColor(0x0D, 0x2E, 0x4A), line=TEAL)
    txb(s, "\u276f  " + q, Inches(5.35), qy + Inches(0.04), Inches(7.3), Inches(0.38),
        size=13, bold=True, color=TEAL)
    qy += Inches(0.46)
    txb(s, "    " + ans, Inches(5.35), qy, Inches(7.3), Inches(0.38),
        size=13, color=GREY, italic=True)
    qy += Inches(0.6)

add_logo(s)


# ── Save ───────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "Summit_Demo_Deck.pptx")
prs.save(out)
print(f"Saved: {out}")
